from langdetect import detect
from langdetect import detect_langs
import os
from pptx import Presentation

print "--Powerpoint Program Started--"

directory = "OUTPUT\\POWERPOINT\\"

PowerpointFilesOutput = open(directory+"Powerpoint_Output.csv", "w")
PowerpointFileErrorLog = open(directory+"Powerpoint_Error_Log.csv", "w")

PowerpointFilesOutput.write("Filename,File Type,Language1,Language2,Language3"+"\n")
PowerpointFileErrorLog.write("Filename,File Type,Error_Message,Error_Output"+"\n")

counter = 1
filetype = "POWERPOINT"

def ppt_processing(filename):    
    prs = Presentation(filename)
    holder = []    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    holder.append(run.text)
                    
    langHolder = '\n'.join(holder)
    return(detect_langs(langHolder))

directoryLocation = "INPUT\\POWERPOINT\\"

for filename in os.listdir(directoryLocation):
    print counter,": ", filename
    fileToPass = str(directoryLocation+filename)
    if filename[0] != "~":
        try:
            #print filename, " - ", ppt_processing(fileToPass)
            langOutput = ppt_processing(fileToPass)
            numLanguage = str(langOutput).count(':'); stopIter1 =  str(langOutput).find(':'); stopIter2 =  str(langOutput).find(':', stopIter1+1); stopIter3 =  str(langOutput).find(':', stopIter2+1)
            if numLanguage == 1:
                #print filename, "-", str(langOutput)[1:stopIter1]
                PowerpointFilesOutput.write(filename+","+filetype+","+str(langOutput)[1:stopIter1]+",,")
            elif numLanguage == 2:
                #print filename, "-", str(langOutput)[1:stopIter1], str(langOutput)[stopIter2-2:stopIter2] 
                PowerpointFilesOutput.write(filename+","+filetype+","+str(langOutput)[1:stopIter1]+","+str(langOutput)[stopIter2-2:stopIter2]+",")
            elif numLanguage == 3:
                #print filename, "-", str(langOutput)[1:stopIter1], str(langOutput)[stopIter2-2:stopIter2], str(langOutput)[stopIter3-2:stopIter3] 
                PowerpointFilesOutput.write(filename+","+filetype+","+str(langOutput)[1:stopIter1]+","+str(langOutput)[stopIter2-2:stopIter2]+","+str(langOutput)[stopIter3-2:stopIter3])
            elif numLanguage > 3:
                #print filename, "-", str(langOutput)[1:stopIter1], str(langOutput)[stopIter2-2:stopIter2], str(langOutput)[stopIter3-2:stopIter3] 
                PowerpointFilesOutput.write(filename+","+filetype+","+str(langOutput)[1:stopIter1]+","+str(langOutput)[stopIter2-2:stopIter2]+","+str(langOutput)[stopIter3-2:stopIter3])
        except Exception as e:
            #print filename, "did not work"
            #print filename, "did not work"
            PowerpointFileErrorLog.write(filename+","+filetype+","+"File Encoding Issue"+","+str(e))
    else:
        continue
    PowerpointFilesOutput.write("\n")
    PowerpointFileErrorLog.write("\n")
    counter +=1
    
PowerpointFilesOutput.close()
PowerpointFileErrorLog.close()

print "--Powerpoint Program Complete--"